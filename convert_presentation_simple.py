"""
Convert LegalNexus Presentation Markdown to PowerPoint and PDF
"""

import re
import sys
from pathlib import Path

# Import required packages (already installed)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_CENTER

def parse_markdown_slides(md_file):
    """Parse markdown file into individual slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide separator (---)
    slides = content.split('\n---\n')

    parsed_slides = []
    for slide in slides:
        slide = slide.strip()
        if not slide or len(slide) < 10:
            continue

        lines = slide.split('\n')

        # Extract title
        title = ""
        content_lines = []

        for line in lines:
            if line.startswith('## Slide'):
                continue
            elif line.startswith('##'):
                title = line.replace('##', '').strip()
            elif line.startswith('#'):
                title = line.replace('#', '').strip()
            elif line.strip():
                content_lines.append(line)

        if title:  # Only add if has title
            parsed_slides.append({
                'title': title,
                'content': '\n'.join(content_lines)
            })

    return parsed_slides

def clean_text(text):
    """Remove markdown formatting"""
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    text = re.sub(r'`([^`]+)`', r'\1', text)
    return text

def create_powerpoint(slides, output_file):
    """Create PowerPoint presentation"""
    print("Creating PowerPoint presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colors
    DARK_SLATE = RGBColor(44, 62, 80)
    GOLD = RGBColor(243, 156, 18)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(245, 245, 245)

    for idx, slide_data in enumerate(slides):
        print(f"  Slide {idx + 1}/{len(slides)}: {slide_data['title'][:40]}...")

        # Blank layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_GRAY if idx == 0 else WHITE

        # Title
        left = Inches(0.5)
        top = Inches(0.4)
        width = Inches(9)
        height = Inches(1.2)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data['title']
        title_p.font.size = Pt(36 if idx == 0 else 32)
        title_p.font.bold = True
        title_p.font.color.rgb = DARK_SLATE
        title_p.alignment = PP_ALIGN.CENTER

        # Content
        left = Inches(0.7)
        top = Inches(1.8)
        width = Inches(8.6)
        height = Inches(5.2)

        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = 1  # Top

        # Parse content
        content_lines = slide_data['content'].split('\n')
        first_para = True

        for line in content_lines:
            line = line.strip()
            if not line:
                continue

            line = clean_text(line)

            # Bullet points
            is_bullet = line.startswith(('-', '*', '•'))
            if is_bullet:
                line = line.lstrip('-*• ').strip()

            # Add paragraph
            if first_para:
                p = content_frame.paragraphs[0]
                first_para = False
            else:
                p = content_frame.add_paragraph()

            p.text = line
            p.font.size = Pt(18 if idx == 0 else 16)
            p.font.color.rgb = DARK_SLATE
            p.level = 1 if is_bullet else 0
            p.space_after = Pt(12)

            # Center align for cover slide
            if idx == 0:
                p.alignment = PP_ALIGN.CENTER

    prs.save(str(output_file))
    print(f"[OK] PowerPoint saved: {output_file}")
    print(f"     Total slides: {len(slides)}")

def create_pdf(slides, output_file):
    """Create PDF document"""
    print("Creating PDF document...")

    doc = SimpleDocTemplate(
        str(output_file),
        pagesize=letter,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=18,
    )

    elements = []
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor='#2c3e50',
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )

    slide_title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading2'],
        fontSize=18,
        textColor='#2c3e50',
        spaceAfter=12,
        fontName='Helvetica-Bold'
    )

    content_style = ParagraphStyle(
        'Content',
        parent=styles['BodyText'],
        fontSize=11,
        textColor='#333333',
        spaceAfter=6,
        leftIndent=20
    )

    # Cover page
    elements.append(Spacer(1, 2*inch))
    elements.append(Paragraph("LegalNexus Enterprise", title_style))
    elements.append(Spacer(1, 0.3*inch))
    elements.append(Paragraph("AI-Powered Legal Practice Management System", content_style))
    elements.append(Spacer(1, 0.5*inch))
    elements.append(Paragraph("Powered by Vicktoria AI - Soul Logic Technology", content_style))
    elements.append(PageBreak())

    # Add slides
    for idx, slide_data in enumerate(slides[1:], 1):  # Skip cover
        print(f"  Adding slide {idx}/{len(slides)-1} to PDF...")

        # Title
        elements.append(Paragraph(slide_data['title'], slide_title_style))
        elements.append(Spacer(1, 0.2*inch))

        # Content
        content_lines = slide_data['content'].split('\n')
        for line in content_lines:
            line = line.strip()
            if not line:
                elements.append(Spacer(1, 0.1*inch))
                continue

            line = clean_text(line)

            # Bullets
            if line.startswith(('-', '*', '•')):
                line = '  • ' + line.lstrip('-*• ').strip()

            # Escape HTML
            line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

            try:
                elements.append(Paragraph(line, content_style))
            except:
                pass  # Skip problematic lines

        elements.append(Spacer(1, 0.3*inch))

        if idx < len(slides) - 1:
            elements.append(PageBreak())

    doc.build(elements)
    print(f"[OK] PDF saved: {output_file}")
    print(f"     Total pages: {len(slides)}")

def main():
    """Main function"""
    print("=" * 60)
    print("LegalNexus Presentation Converter")
    print("=" * 60)

    md_file = Path("LEGALNEXUS_PRESENTATION.md")

    if not md_file.exists():
        print(f"ERROR: {md_file} not found!")
        sys.exit(1)

    print(f"Reading: {md_file}")

    # Parse
    slides = parse_markdown_slides(md_file)
    print(f"Parsed {len(slides)} slides\n")

    # Output directory
    output_dir = Path("presentation_output")
    output_dir.mkdir(exist_ok=True)

    # Convert
    pptx_file = output_dir / "LegalNexus_Presentation.pptx"
    pdf_file = output_dir / "LegalNexus_Presentation.pdf"

    try:
        create_powerpoint(slides, pptx_file)
    except Exception as e:
        print(f"PowerPoint error: {e}")

    print()

    try:
        create_pdf(slides, pdf_file)
    except Exception as e:
        print(f"PDF error: {e}")

    print("\n" + "=" * 60)
    print("CONVERSION COMPLETE!")
    print("=" * 60)
    print(f"\nOutput directory: {output_dir.absolute()}")
    print(f"  - PowerPoint: {pptx_file.name}")
    print(f"  - PDF: {pdf_file.name}")
    print("\nReady to present!")

if __name__ == "__main__":
    main()
